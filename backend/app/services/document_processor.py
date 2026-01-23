"""
Document Processor Service - Multi-format document processing
Supports: PDF, Word, Excel, PowerPoint, Text files
"""

import logging
from typing import List, Dict, Any, Optional
from pathlib import Path
import hashlib
import os

# PDF processing
import PyPDF2
import pdfplumber

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Service for processing multiple document formats"""
    
    SUPPORTED_FORMATS = {
        'pdf': ['.pdf'],
        'word': ['.docx', '.doc'],
        'excel': ['.xlsx', '.xls', '.csv'],
        'powerpoint': ['.pptx', '.ppt'],
        'text': ['.txt', '.md']
    }
    
    def __init__(self):
        """Initialize document processor"""
        self.pdf_available = True
        self.word_available = self._check_word_support()
        self.excel_available = self._check_excel_support()
        self.ppt_available = self._check_ppt_support()
        
        logger.info(f"Document processor initialized")
        logger.info(f"  PDF: {self.pdf_available}")
        logger.info(f"  Word: {self.word_available}")
        logger.info(f"  Excel: {self.excel_available}")
        logger.info(f"  PowerPoint: {self.ppt_available}")
    
    def _check_word_support(self) -> bool:
        """Check if Word processing is available"""
        try:
            import docx
            return True
        except ImportError:
            logger.warning("python-docx not installed - Word support disabled")
            return False
    
    def _check_excel_support(self) -> bool:
        """Check if Excel processing is available"""
        try:
            import openpyxl
            import pandas
            return True
        except ImportError:
            logger.warning("openpyxl/pandas not installed - Excel support disabled")
            return False
    
    def _check_ppt_support(self) -> bool:
        """Check if PowerPoint processing is available"""
        try:
            import pptx
            return True
        except ImportError:
            logger.warning("python-pptx not installed - PowerPoint support disabled")
            return False
    
    def get_file_type(self, filename: str) -> Optional[str]:
        """
        Determine file type from extension.
        
        Returns:
            File type ('pdf', 'word', 'excel', 'powerpoint', 'text') or None
        """
        ext = Path(filename).suffix.lower()
        
        for file_type, extensions in self.SUPPORTED_FORMATS.items():
            if ext in extensions:
                return file_type
        
        return None
    
    def is_supported(self, filename: str) -> bool:
        """Check if file format is supported"""
        file_type = self.get_file_type(filename)
        
        if not file_type:
            return False
        
        # Check if required library is available
        if file_type == 'word' and not self.word_available:
            return False
        if file_type == 'excel' and not self.excel_available:
            return False
        if file_type == 'powerpoint' and not self.ppt_available:
            return False
        
        return True
    
    def extract_text(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text from any supported document format.
        
        Returns:
            List of dicts with 'page'/'sheet'/'slide' and 'text' keys
        """
        filename = Path(file_path).name
        file_type = self.get_file_type(filename)
        
        if file_type == 'pdf':
            return self._extract_pdf(file_path)
        elif file_type == 'word':
            return self._extract_word(file_path)
        elif file_type == 'excel':
            return self._extract_excel(file_path)
        elif file_type == 'powerpoint':
            return self._extract_ppt(file_path)
        elif file_type == 'text':
            return self._extract_text(file_path)
        else:
            raise ValueError(f"Unsupported file type: {filename}")
    
    def _extract_pdf(self, pdf_path: str) -> List[Dict[str, Any]]:
        """Extract text from PDF"""
        pages = []
        
        try:
            # Check if file exists and get size
            if not os.path.exists(pdf_path):
                raise FileNotFoundError(f"PDF file not found: {pdf_path}")
            
            file_size = os.path.getsize(pdf_path)
            logger.debug(f"Opening PDF: {pdf_path} (size: {file_size} bytes)")
            
            with pdfplumber.open(pdf_path) as pdf:
                total_pages = len(pdf.pages)
                logger.debug(f"PDF has {total_pages} total pages")
                
                for page_num, page in enumerate(pdf.pages, start=1):
                    try:
                        text = page.extract_text()
                        if text and text.strip():
                            pages.append({
                                'page': page_num,
                                'text': text.strip()
                            })
                            logger.debug(f"Page {page_num}: extracted {len(text)} characters")
                        else:
                            logger.warning(f"Page {page_num}: no text extracted (might be image-only)")
                    except Exception as page_error:
                        logger.error(f"Error extracting page {page_num}: {page_error}")
                        # Continue to next page instead of failing completely
                        continue
            
            logger.info(f"Extracted {len(pages)} pages with text from {total_pages} total pages")
            
            if len(pages) == 0:
                error_msg = (
                    f"No text extracted from PDF. This is likely a scanned/image-based PDF without a text layer. "
                    f"To process this document, you'll need OCR (Optical Character Recognition). "
                    f"Consider converting the PDF to a text-based format first, or contact support for OCR capabilities."
                )
                logger.error(error_msg)
                raise ValueError("Image-only PDF detected: No extractable text found. OCR required for scanned documents.")
            
            return pages
            
        except Exception as e:
            logger.error(f"Error extracting PDF from {pdf_path}: {type(e).__name__}: {e}")
            raise
    
    def _extract_word(self, docx_path: str) -> List[Dict[str, Any]]:
        """Extract text from Word document"""
        if not self.word_available:
            raise ImportError("python-docx not installed. Install with: pip install python-docx")
        
        try:
            import docx
            
            doc = docx.Document(docx_path)
            pages = []
            
            # Word doesn't have pages, so we treat each paragraph as a section
            current_page = 1
            page_text = []
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    page_text.append(text)
                    
                    # Create a "page" every 10 paragraphs or at page breaks
                    if len(page_text) >= 10:
                        pages.append({
                            'page': current_page,
                            'text': '\n'.join(page_text)
                        })
                        page_text = []
                        current_page += 1
            
            # Add remaining text
            if page_text:
                pages.append({
                    'page': current_page,
                    'text': '\n'.join(page_text)
                })
            
            logger.info(f"Extracted {len(pages)} sections from Word document")
            return pages
            
        except Exception as e:
            logger.error(f"Error extracting Word document: {e}")
            raise
    
    def _extract_excel(self, excel_path: str) -> List[Dict[str, Any]]:
        """Extract text from Excel file"""
        if not self.excel_available:
            raise ImportError("openpyxl/pandas not installed. Install with: pip install openpyxl pandas")
        
        try:
            import pandas as pd
            
            # Read all sheets
            excel_file = pd.ExcelFile(excel_path)
            sheets = []
            
            for sheet_num, sheet_name in enumerate(excel_file.sheet_names, start=1):
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # Convert dataframe to text
                text_parts = [f"Sheet: {sheet_name}\n"]
                
                # Add column headers
                text_parts.append("Columns: " + ", ".join(df.columns.astype(str)))
                
                # Add rows (limit to prevent huge text)
                for idx, row in df.head(100).iterrows():
                    row_text = " | ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                    if row_text:
                        text_parts.append(row_text)
                
                sheets.append({
                    'page': sheet_num,  # Use 'page' for consistency
                    'sheet': sheet_name,
                    'text': '\n'.join(text_parts)
                })
            
            logger.info(f"Extracted {len(sheets)} sheets from Excel")
            return sheets
            
        except Exception as e:
            logger.error(f"Error extracting Excel: {e}")
            raise
    
    def _extract_ppt(self, ppt_path: str) -> List[Dict[str, Any]]:
        """Extract text from PowerPoint"""
        if not self.ppt_available:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")
        
        try:
            from pptx import Presentation
            
            prs = Presentation(ppt_path)
            slides = []
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                text_parts = []
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            text_parts.append(text)
                
                if text_parts:
                    slides.append({
                        'page': slide_num,  # Use 'page' for consistency
                        'slide': slide_num,
                        'text': '\n'.join(text_parts)
                    })
            
            logger.info(f"Extracted {len(slides)} slides from PowerPoint")
            return slides
            
        except Exception as e:
            logger.error(f"Error extracting PowerPoint: {e}")
            raise
    
    def _extract_text(self, text_path: str) -> List[Dict[str, Any]]:
        """Extract text from plain text file"""
        try:
            with open(text_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            # Split into chunks (treat as single page)
            return [{
                'page': 1,
                'text': text.strip()
            }]
            
        except Exception as e:
            logger.error(f"Error reading text file: {e}")
            raise
    
    def chunk_text(
        self, 
        text: str, 
        chunk_size: int = 1000, 
        overlap: int = 200
    ) -> List[str]:
        """
        Split text into overlapping chunks.
        """
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            end = start + chunk_size
            chunk = text[start:end]
            
            # Try to break at sentence boundary
            if end < text_length:
                last_period = chunk.rfind('.')
                last_newline = chunk.rfind('\n')
                break_point = max(last_period, last_newline)
                
                if break_point > chunk_size * 0.5:
                    chunk = chunk[:break_point + 1]
                    end = start + break_point + 1
            
            chunks.append(chunk.strip())
            start = end - overlap
        
        return chunks
    
    def process_document(
        self,
        file_path: str,
        doc_metadata: Dict[str, Any],
        chunk_size: int = 1000,
        overlap: int = 200
    ) -> List[Dict[str, Any]]:
        """
        Process any supported document into chunks with metadata.
        """
        # Extract text
        pages = self.extract_text(file_path)
        
        # Process each page/sheet/slide
        all_chunks = []
        chunk_id = 0
        doc_id = doc_metadata.get('doc_id')
        
        for page_data in pages:
            page_num = page_data['page']
            page_text = page_data['text']
            
            # Chunk the text
            chunks = self.chunk_text(page_text, chunk_size, overlap)
            
            # Add metadata to each chunk
            for chunk_text in chunks:
                chunk_id += 1
                # Generate unique ID using hash of doc_id + chunk_id
                unique_id = hashlib.md5(f"{doc_id}_{chunk_id}".encode()).hexdigest()
                chunk_metadata = {
                    'id': unique_id,
                    'text': chunk_text,
                    'page': page_num,
                    'doc_id': doc_id,
                    'title': doc_metadata.get('title'),
                    'program': doc_metadata.get('program'),
                    'department': doc_metadata.get('department'),
                    'semester': doc_metadata.get('semester'),
                    'category': doc_metadata.get('category'),
                    'version': doc_metadata.get('version', 1)
                }
                all_chunks.append(chunk_metadata)
        
        logger.debug(f"Created {len(all_chunks)} chunks from {len(pages)} pages/sections")
        return all_chunks
    
    def generate_doc_id(self, filename: str, metadata: Dict[str, Any]) -> str:
        """Generate unique document ID"""
        unique_string = f"{filename}_{metadata.get('program')}_{metadata.get('department')}_{metadata.get('semester')}_{metadata.get('category')}"
        return hashlib.md5(unique_string.encode()).hexdigest()[:16]


# Singleton instance
_document_processor = None

def get_document_processor() -> DocumentProcessor:
    """Get or create document processor singleton"""
    global _document_processor
    if _document_processor is None:
        _document_processor = DocumentProcessor()
    return _document_processor
